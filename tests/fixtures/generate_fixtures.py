"""
Generates the small .pptx fixtures under tests/fixtures/.

Each fixture is built to trip exactly one check in scripts/slop_check.py, so
tests/test_slop_check.py can assert "this input produces that finding" instead
of eyeballing example decks. Re-run this after changing what a fixture needs
to contain:

    python tests/fixtures/generate_fixtures.py

Requires python-pptx and Pillow (dev-only — not in scripts/requirements.txt,
since fixtures are committed as binaries and regenerated only when edited).
"""

import io
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image

HERE = Path(__file__).resolve().parent

TITLE_AND_CONTENT = 1
TITLE_SLIDE = 0
BLANK = 6


def new_deck():
    return Presentation()


def set_run_size(paragraph, size_pt):
    for run in paragraph.runs:
        run.font.size = Pt(size_pt)


def add_slide(prs, title=None, body_lines=None, notes=None,
              title_size=None, body_size=None, layout=TITLE_AND_CONTENT):
    slide = prs.slides.add_slide(prs.slide_layouts[layout])

    if title is not None and slide.shapes.title is not None:
        slide.shapes.title.text = title
        if title_size is not None:
            set_run_size(slide.shapes.title.text_frame.paragraphs[0], title_size)

    if body_lines:
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = body_lines[0]
        if body_size is not None:
            set_run_size(tf.paragraphs[0], body_size)
        for line in body_lines[1:]:
            p = tf.add_paragraph()
            p.text = line
            if body_size is not None:
                set_run_size(p, body_size)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def make_image(path, size, fmt, color=(15, 42, 74)):
    Image.new("RGB", size, color).save(path, fmt)


def save(prs, name):
    out = HERE / name
    prs.save(out)
    print(f"wrote {out.relative_to(HERE.parent.parent)}")


def build_clean():
    """Baseline: should clear every check (used as the negative case)."""
    prs = new_deck()
    logo_path = HERE / "_tmp_logo_clean.png"
    make_image(logo_path, (280, 200), "PNG")  # 1.4:1 aspect

    s1 = add_slide(prs, title="Q3 Operations Review", title_size=32, layout=TITLE_SLIDE)
    s1.shapes.add_picture(str(logo_path), Inches(11.2), Inches(0.4), width=Inches(1.4), height=Inches(1.0))

    add_slide(prs, title="Ticket Volume Fell 18% After Rollout",
              body_lines=["Password-reset tickets dropped from 420/week to 344/week.",
                          "Your Company — Confidential"],
              title_size=28, body_size=18)
    add_slide(prs, title="Next Steps for Q4",
              body_lines=["Expand rollout to the billing team.",
                          "Your Company — Confidential"],
              title_size=28, body_size=18)
    save(prs, "clean.pptx")
    logo_path.unlink()


def build_banned_phrase():
    prs = new_deck()
    add_slide(prs, title="Platform Overview",
              body_lines=["Our seamless integration lets teams move faster."],
              title_size=28, body_size=18)
    save(prs, "banned_phrase.pptx")


def build_banned_phrase_notes():
    prs = new_deck()
    add_slide(prs, title="Platform Overview",
              body_lines=["Adoption reached 42% among pilot users."],
              notes="Remember to say this is a real game-changer live.",
              title_size=28, body_size=18)
    save(prs, "banned_phrase_notes.pptx")


def build_empty_slide():
    prs = new_deck()
    add_slide(prs, title="Overview", body_lines=["Something concrete happens here."],
              title_size=28, body_size=18)
    prs.slides.add_slide(prs.slide_layouts[BLANK])  # no text at all
    save(prs, "empty_slide.pptx")


def build_content_free_slide():
    prs = new_deck()
    add_slide(prs, title="Overview", body_lines=["Something concrete happens here."],
              title_size=28, body_size=18)
    add_slide(prs, title="Thank You", title_size=28)
    save(prs, "content_free_slide.pptx")


def build_title_only_slide():
    prs = new_deck()
    add_slide(prs, title="Overview", body_lines=["Something concrete happens here."],
              title_size=28, body_size=18)
    add_slide(prs, title="Pricing", title_size=28)
    save(prs, "title_only_slide.pptx")


def build_low_specificity():
    prs = new_deck()
    add_slide(prs, title="Overview",
              body_lines=["our strategy centers on growth and long term value creation"],
              title_size=28, body_size=18)
    save(prs, "low_specificity.pptx")


def build_repeated_openings():
    # check_repeated_openings keys on the first 3 letter-tokens of the title,
    # so both need to match there and diverge after.
    prs = new_deck()
    add_slide(prs, title="Project Update Phase Rollout", body_lines=["Shipped to 3 pilot teams."],
              title_size=28, body_size=18)
    add_slide(prs, title="Project Update Phase Billing", body_lines=["Targets the billing team next."],
              title_size=28, body_size=18)
    save(prs, "repeated_openings.pptx")


def build_exclamation_mark():
    prs = new_deck()
    add_slide(prs, title="Overview",
              body_lines=["Adoption hit 42% among pilot users!"],
              title_size=28, body_size=18)
    save(prs, "exclamation_mark.pptx")


def build_rhetorical_question():
    prs = new_deck()
    add_slide(prs, title="What If We Moved Faster?",
              body_lines=["Cycle time dropped from 9 days to 4 days."],
              title_size=28, body_size=18)
    save(prs, "rhetorical_question.pptx")


def build_slide_cap_exceeded():
    prs = new_deck()
    for i in range(1, 13):  # 12 slides, default cap is 10
        add_slide(prs, title=f"Topic {i}", body_lines=[f"Point {i} with some detail."],
                  title_size=28, body_size=18)
    save(prs, "slide_cap_exceeded.pptx")


def build_font_below_minimum():
    prs = new_deck()
    add_slide(prs, title="Overview",
              body_lines=["This body text is set below the brand minimum."],
              title_size=28, body_size=10)  # min_body default is 14
    save(prs, "font_below_minimum.pptx")


def build_font_size_unresolved():
    prs = new_deck()
    # No explicit sizes set anywhere -> runs inherit from the layout/master.
    add_slide(prs, title="Overview", body_lines=["Sizes here are never set explicitly."])
    save(prs, "font_size_unresolved.pptx")


def build_missing_footer():
    prs = new_deck()
    add_slide(prs, title="Title Slide", title_size=32, layout=TITLE_SLIDE)
    add_slide(prs, title="Adoption Update",
              body_lines=["Adoption reached 42% among pilot users."],
              title_size=28, body_size=18)  # no footer text anywhere
    save(prs, "missing_footer.pptx")


def build_logo_missing():
    prs = new_deck()
    add_slide(prs, title="Title Slide", title_size=32, layout=TITLE_SLIDE)  # no picture
    add_slide(prs, title="Overview", body_lines=["Adoption reached 42% among pilot users."],
              title_size=28, body_size=18)
    save(prs, "logo_missing.pptx")


def build_logo_format():
    prs = new_deck()
    logo_path = HERE / "_tmp_logo_format.jpg"
    make_image(logo_path, (280, 200), "JPEG")
    s1 = add_slide(prs, title="Title Slide", title_size=32, layout=TITLE_SLIDE)
    s1.shapes.add_picture(str(logo_path), Inches(11.2), Inches(0.4), width=Inches(1.4), height=Inches(1.0))
    add_slide(prs, title="Overview", body_lines=["Adoption reached 42% among pilot users."],
              title_size=28, body_size=18)
    save(prs, "logo_format.pptx")
    logo_path.unlink()


def build_logo_stretched():
    prs = new_deck()
    logo_path = HERE / "_tmp_logo_stretch.png"
    make_image(logo_path, (400, 100), "PNG")  # native 4:1
    s1 = add_slide(prs, title="Title Slide", title_size=32, layout=TITLE_SLIDE)
    s1.shapes.add_picture(str(logo_path), Inches(11.0), Inches(0.4), width=Inches(2.0), height=Inches(2.0))  # square placement
    add_slide(prs, title="Overview", body_lines=["Adoption reached 42% among pilot users."],
              title_size=28, body_size=18)
    save(prs, "logo_stretched.pptx")
    logo_path.unlink()


def build_logo_too_small():
    prs = new_deck()
    logo_path = HERE / "_tmp_logo_small.png"
    make_image(logo_path, (400, 100), "PNG")  # native 4:1
    s1 = add_slide(prs, title="Title Slide", title_size=32, layout=TITLE_SLIDE)
    s1.shapes.add_picture(str(logo_path), Inches(11.8), Inches(0.4), width=Inches(0.5), height=Inches(0.125))
    add_slide(prs, title="Overview", body_lines=["Adoption reached 42% among pilot users."],
              title_size=28, body_size=18)
    save(prs, "logo_too_small.pptx")
    logo_path.unlink()


def build_logo_size_off_spec():
    prs = new_deck()
    logo_path = HERE / "_tmp_logo_offspec.png"
    make_image(logo_path, (400, 100), "PNG")  # native 4:1
    s1 = add_slide(prs, title="Title Slide", title_size=32, layout=TITLE_SLIDE)
    s1.shapes.add_picture(str(logo_path), Inches(11.4), Inches(0.4), width=Inches(1.0), height=Inches(0.25))
    add_slide(prs, title="Overview", body_lines=["Adoption reached 42% among pilot users."],
              title_size=28, body_size=18)
    save(prs, "logo_size_off_spec.pptx")
    logo_path.unlink()


if __name__ == "__main__":
    build_clean()
    build_banned_phrase()
    build_banned_phrase_notes()
    build_empty_slide()
    build_content_free_slide()
    build_title_only_slide()
    build_low_specificity()
    build_repeated_openings()
    build_exclamation_mark()
    build_rhetorical_question()
    build_slide_cap_exceeded()
    build_font_below_minimum()
    build_font_size_unresolved()
    build_missing_footer()
    build_logo_missing()
    build_logo_format()
    build_logo_stretched()
    build_logo_too_small()
    build_logo_size_off_spec()
