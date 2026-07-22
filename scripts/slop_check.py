import argparse
import json
import re
import sys
from collections import defaultdict
from pathlib import Path

TOOL_ERR = "SLOP_CHECK_TOOL_ERROR:"

try:
    from pptx import Presentation
except ImportError:
    print(f"{TOOL_ERR} python-pptx not installed (pip install python-pptx)", file=sys.stderr)
    sys.exit(2)

try:
    import yaml
except ImportError:
    yaml = None

EMU_PER_INCH = 914400

#default branding styles
DEFAULT_BRAND = {
    "type_scale": {"min_title": 24, "min_body": 14},
    "slide_caps": {"default": 10},
    "footer": {"client_facing_text": "Company-X — Confidential"},
    "logo": {"width_in": 1.4, "min_width_in": 0.9},
}

#load branding styles from brand.yaml
def load_brand(explicit=None):
    """brand.yaml lives at the skill root, one level above scripts/."""
    if explicit:
        path = Path(explicit)
    else:
        path = Path(__file__).resolve().parent.parent / "brand.yaml"
    if yaml is None or not path.exists():
        return DEFAULT_BRAND, None
    try:
        with open(path, encoding="utf-8") as fh:
            data = yaml.safe_load(fh) or {}
    except Exception as exc:
        return DEFAULT_BRAND, f"could not parse {path}: {exc}"
    merged = {**DEFAULT_BRAND, **data}
    for key in DEFAULT_BRAND:
        if isinstance(DEFAULT_BRAND[key], dict):
            merged[key] = {**DEFAULT_BRAND[key], **(data.get(key) or {})}
    return merged, None

# anti-slop phrases
BANNED_PHRASES = [
    r"fast[\s-]?paced world",
    r"unlock the power",
    r"seamless(ly)?",
    r"robust solution",
    r"cutting[\s-]?edge",
    r"in today'?s",
    r"let'?s dive in",
    r"dive into",
    r"game[\s-]?changer",
    r"at the end of the day",
    r"synerg(y|ies|istic)",
    r"\bleverag(e|es|ing)\b(?!\s+(?:ratio|point|effect))",
    r"paradigm shift",
    r"take .* to the next level",
    r"holistic approach",
    r"best[\s-]?in[\s-]?class",
    r"world[\s-]?class",
    r"empower(ing|ed|s)?",
    r"revolutioniz(e|ing|ed)",
    r"transformative",
    r"innovative solutions?",
]
BANNED_RE = re.compile("|".join(BANNED_PHRASES), re.IGNORECASE)
ABSTRACT_BUZZWORDS = {
    "strategy", "growth", "value", "innovation", "excellence", "solutions",
    "synergy", "efficiency", "optimization", "transformation", "impact",
    "engagement", "alignment", "vision", "mission", "empowerment",
}

CONTENT_FREE_TITLES = {
    "thank you", "thanks", "questions", "questions?", "q&a", "the end",
}

CAP_STOPWORDS = {
    "The", "A", "An", "We", "Our", "Your", "This", "That", "These", "Those",
    "It", "In", "On", "For", "And", "But", "With", "To", "Of", "By", "At",
    "As", "If", "When", "How", "What", "Why", "Is", "Are", "Be", "Will",
    "Can", "All", "Each", "Every", "New", "More", "Most", "Less", "Both",
} | {w.capitalize() for w in ABSTRACT_BUZZWORDS}

NUMBER_PATTERN = re.compile(
    r"(?:\$\s?\d+(?:,\d{3})*(?:\.\d+)?\s?[kmb]?\b"      
    r"|\b\d{1,3}(?:,\d{3})+(?:\.\d+)?\b"               
    r"|\b\d+(?:\.\d+)?\s?%"                             
    r"|\b\d+(?:\.\d+)?\b)",                           
    re.IGNORECASE,
)
YEAR_RE = re.compile(r"^(?:19|20)\d{2}$")

HIGH, MEDIUM, LOW, INFO = "high", "medium", "low", "info"
SEVERITY_ORDER = {HIGH: 0, MEDIUM: 1, LOW: 2, INFO: 3}
BLOCKING = (HIGH, MEDIUM)

def finding(slide, check, detail, severity, scope="slide"):
    return {"slide": slide, "check": check, "detail": detail,
            "severity": severity, "scope": scope}

#extraction phase
class Slide:
    def __init__(self, number):
        self.number = number
        self.title = ""
        self.body_lines = []
        self.notes = ""
        self.runs = []       
        self.pictures = []   

    @property
    def visible_text(self):
        return "\n".join(([self.title] if self.title else []) + self.body_lines)

    @property
    def body_text(self):
        return "\n".join(self.body_lines)
    
def extract(pptx_path):
    prs = Presentation(pptx_path)
    slides = []
    for idx, raw in enumerate(prs.slides, start=1):
        s = Slide(idx)

        title_shape = None
        try:
            title_shape = raw.shapes.title
        except Exception:
            pass

        for shape in raw.shapes:
            try:
                img = shape.image
                s.pictures.append((shape.width, shape.height,
                                   img.size[0], img.size[1], (img.ext or "").lower()))
                continue
            except Exception:
                pass

            if not getattr(shape, "has_text_frame", False):
                continue

            is_title = title_shape is not None and shape is title_shape
            for para in shape.text_frame.paragraphs:
                line = "".join(r.text for r in para.runs).strip()
                for r in para.runs:
                    size = r.font.size.pt if r.font.size is not None else None
                    if r.text.strip():
                        s.runs.append((r.text.strip(), size, is_title))
                if not line:
                    continue
                if is_title and not s.title:
                    s.title = line
                elif is_title:
                    s.body_lines.append(line)
                else:
                    s.body_lines.append(line)

        if not s.title and s.body_lines:
            s.title = s.body_lines.pop(0)

        if raw.has_notes_slide:
            s.notes = (raw.notes_slide.notes_text_frame.text or "").strip()

        slides.append(s)
    return slides, prs

#text checks
def check_banned_phrases(slides):
    out = []
    for s in slides:
        for m in BANNED_RE.finditer(s.visible_text):
            out.append(finding(s.number, "banned_phrase",
                               f"phrase matching /{m.group(0)}/", HIGH))
        for m in BANNED_RE.finditer(s.notes):
            out.append(finding(s.number, "banned_phrase",
                               f"in speaker notes: /{m.group(0)}/", LOW, scope="notes"))
    return out


def check_content_free_slides(slides):
    out = []
    for s in slides:
        if not s.title and not s.body_lines:
            out.append(finding(s.number, "empty_slide",
                               "no extractable text at all", HIGH))
            continue
        norm = s.title.lower().strip("*#: ")
        if norm in CONTENT_FREE_TITLES and not s.body_lines:
            out.append(finding(s.number, "content_free_slide",
                               f"filler slide ('{s.title}') with no content", MEDIUM))
        elif not s.body_lines and s.number != 1:
            # Was previously missed: any title-only slide, not just the whitelist.
            out.append(finding(s.number, "title_only_slide",
                               f"title '{s.title}' with no body content", MEDIUM))
    return out


def _has_named_entity(text):
    for line in text.splitlines():
        tokens = re.findall(r"\b[A-Za-z][A-Za-z0-9&.\-]*\b", line)
        if not tokens:
            continue
        if any(len(t) >= 2 and t.isupper() for t in tokens):
            return True
        for t in tokens[1:]:
            if t[0].isupper() and t not in CAP_STOPWORDS:
                return True
    return False


def check_specificity(slides):
    out = []
    for s in slides:
        body = s.body_text
        words = re.findall(r"[a-zA-Z]+", body.lower())
        if not words:
            continue
        buzz = sum(1 for w in words if w in ABSTRACT_BUZZWORDS)
        if buzz >= 2 and not NUMBER_PATTERN.search(body) and not _has_named_entity(body):
            out.append(finding(s.number, "low_specificity",
                               f"{buzz} abstract buzzwords, no number or named entity", MEDIUM))
    return out


def check_numeric_claims(slides):
    out = []
    for s in slides:
        seen = set()
        for m in NUMBER_PATTERN.finditer(s.visible_text):
            tok = m.group(0).strip()
            if YEAR_RE.match(tok) or tok in seen:
                continue
            seen.add(tok)
            out.append(finding(s.number, "numeric_claim_needs_source_check",
                               f"'{tok}' — confirm this traces to the user's source material", INFO))
    return out


def check_repeated_openings(slides):
    out = []
    openings = defaultdict(list)
    for s in slides:
        if not s.title:
            continue
        words = re.findall(r"[a-zA-Z]+", s.title.lower())[:3]
        if len(words) < 2:
            continue
        openings[" ".join(words)].append(s.number)
    for key, nums in openings.items():
        if len(nums) > 1:
            out.append(finding(",".join(map(str, nums)), "repeated_title_structure",
                               f"slides {nums} all open with '{key}...' — vary structure", LOW))
    return out


def check_voice_rules(slides):
    out = []
    for s in slides:
        if "!" in s.visible_text:
            out.append(finding(s.number, "exclamation_mark",
                               "brand voice forbids exclamation marks", MEDIUM))
        title = s.title.rstrip()
        if title.endswith("?") and title.lower().strip("*#: ") not in CONTENT_FREE_TITLES:
            out.append(finding(s.number, "rhetorical_question_title",
                               f"title '{title}' is a question — brand voice forbids this", MEDIUM))
    return out

#brand checks
def check_slide_cap(slides, brand, deck_type):
    caps = brand.get("slide_caps", {})
    cap = caps.get(deck_type, caps.get("default"))
    if not cap:
        return []
    if len(slides) > cap:
        return [finding("-", "slide_cap_exceeded",
                        f"{len(slides)} slides exceeds the cap of {cap} for '{deck_type}' — "
                        f"cut detail, split to appendix, or confirm the overage with the user",
                        HIGH)]
    return []


def check_font_sizes(slides, brand):
    ts = brand.get("type_scale", {})
    min_title = ts.get("min_title", 24)
    min_body = ts.get("min_body", 14)
    out, unresolved = [], 0
    for s in slides:
        for text, size, is_title in s.runs:
            if size is None:
                unresolved += 1
                continue
            floor = min_title if is_title else min_body
            if size < floor:
                out.append(finding(s.number, "font_below_minimum",
                                   f"{size:g}pt on '{text[:40]}' is below the "
                                   f"{'title' if is_title else 'body'} minimum of {floor}pt",
                                   HIGH))
    if unresolved:
        out.append(finding("-", "font_size_unresolved",
                           f"{unresolved} runs inherit size from the layout/master — "
                           f"not verifiable here, confirm visually", INFO))
    return out


def check_footer(slides, brand, client_facing):
    if not client_facing:
        return []
    text = (brand.get("footer", {}) or {}).get("client_facing_text", "")
    if not text:
        return []
    needle = text.replace("—", "-").lower()
    out = []
    for s in slides:
        if s.number == 1:
            continue
        haystack = s.visible_text.replace("—", "-").lower()
        if needle not in haystack:
            out.append(finding(s.number, "missing_footer",
                               f"client-facing deck is missing '{text}'", HIGH))
    return out


def check_logo(slides, brand):
    cfg = brand.get("logo", {}) or {}
    out = []
    title_slide = slides[0] if slides else None
    if title_slide is None:
        return out
    if not title_slide.pictures:
        return [finding(1, "logo_missing", "no image found on the title slide", HIGH)]

    for w_emu, h_emu, nat_w, nat_h, ext in title_slide.pictures:
        if ext in ("jpg", "jpeg"):
            out.append(finding(1, "logo_format",
                               "logo is JPEG — no transparency, will render a white box "
                               "on colored fills; use PNG", HIGH))
        if nat_w and nat_h and h_emu:
            native = nat_w / nat_h
            placed = w_emu / h_emu
            if abs(placed - native) / native > 0.02:
                out.append(finding(1, "logo_stretched",
                                   f"placed aspect {placed:.3f} vs native {native:.3f} "
                                   f"— scale proportionally", HIGH))
        width_in = w_emu / EMU_PER_INCH
        min_w = cfg.get("min_width_in")
        target = cfg.get("width_in")
        if min_w and width_in < min_w:
            out.append(finding(1, "logo_too_small",
                               f"{width_in:.2f}in is below the minimum of {min_w}in", MEDIUM))
        elif target and abs(width_in - target) / target > 0.15:
            out.append(finding(1, "logo_size_off_spec",
                               f"{width_in:.2f}in vs spec {target}in", LOW))
    return out

#running antislop checks
def run(pptx_path, deck_type, client_facing, brand_path, strict, as_json):
    brand, brand_err = load_brand(brand_path)

    try:
        slides, _ = extract(pptx_path)
    except Exception as exc:
        print(f"{TOOL_ERR} could not read {pptx_path}: {exc}", file=sys.stderr)
        return 2
    if not slides:
        print(f"{TOOL_ERR} no slides found in {pptx_path}", file=sys.stderr)
        return 2

    findings = []
    if brand_err:
        findings.append(finding("-", "brand_config", brand_err + " — using defaults", INFO))

    findings += check_banned_phrases(slides)
    findings += check_content_free_slides(slides)
    findings += check_specificity(slides)
    findings += check_numeric_claims(slides)
    findings += check_repeated_openings(slides)
    findings += check_voice_rules(slides)
    findings += check_slide_cap(slides, brand, deck_type)
    findings += check_font_sizes(slides, brand)
    findings += check_footer(slides, brand, client_facing)
    findings += check_logo(slides, brand)

    findings.sort(key=lambda f: (SEVERITY_ORDER.get(f["severity"], 9), str(f["slide"])))
    blocking = [f for f in findings if f["severity"] in BLOCKING]

    if as_json:
        print(json.dumps({"file": str(pptx_path), "slides": len(slides),
                          "deck_type": deck_type, "blocking": len(blocking),
                          "findings": findings}, indent=2))
    elif not findings:
        print(f"No findings across {len(slides)} slides.")
    else:
        print(f"{pptx_path}: {len(findings)} findings across {len(slides)} slides\n")
        for f in findings:
            scope = "" if f["scope"] == "slide" else f" ({f['scope']})"
            print(f"  [{f['severity'].upper():6}] slide {str(f['slide']):>5}  "
                  f"{f['check']:32}{scope} {f['detail']}")
        print()
        if blocking:
            print(f"-> {len(blocking)} high/medium findings must be fixed before delivery.")
        else:
            print("-> Only low/info findings — review at your discretion.")

    if strict:
        return 1 if findings else 0
    return 1 if blocking else 0


if __name__ == "__main__":
    p = argparse.ArgumentParser(description="Slop and brand compliance check for a .pptx deck")
    p.add_argument("pptx", help="Path to the .pptx file")
    p.add_argument("--deck-type", default="default",
                   help="Key from slide_caps in brand.yaml (internal_update, "
                        "client_proposal, exec_summary). Default: default")
    p.add_argument("--client-facing", dest="client_facing", action="store_true", default=None,
                   help="Force footer checks on")
    p.add_argument("--internal", dest="client_facing", action="store_false",
                   help="Force footer checks off")
    p.add_argument("--brand", help="Override path to brand.yaml")
    p.add_argument("--strict", action="store_true",
                   help="Exit non-zero on ANY finding, including low/info")
    p.add_argument("--json", dest="as_json", action="store_true",
                   help="Machine-readable output")
    args = p.parse_args()

    cf = args.client_facing
    if cf is None:
        cf = args.deck_type == "client_proposal"

    sys.exit(run(args.pptx, args.deck_type, cf, args.brand, args.strict, args.as_json))