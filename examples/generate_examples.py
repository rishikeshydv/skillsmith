"""
Builds the two example decks from examples/source-brief.md — the same brief,
run through two different processes.

  without-skill.pptx: generic model output with no brand config, no cap, no
  voice rules — the deck you get when you just ask for slides.

  with-skill.pptx: same brief, built against brand.example.yaml following
  SKILL.md's archetypes and caps. Should clear scripts/slop_check.py clean.

Dev-only tool (needs Pillow in addition to scripts/requirements.txt).
Re-run after editing either deck:

    python examples/generate_examples.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from PIL import Image, ImageDraw

HERE = Path(__file__).resolve().parent

TITLE_SLIDE = 0
TITLE_AND_CONTENT = 1
BLANK = 6

PRIMARY = (0x0F, 0x2A, 0x4A)
ACCENT = (0xFF, 0x6B, 0x6B)
FOOTER = "Your Company — Confidential"


def set_run_size(paragraph, size_pt):
    for run in paragraph.runs:
        run.font.size = Pt(size_pt)


def add_slide(prs, title=None, body_lines=None, title_size=32, body_size=18, layout=TITLE_AND_CONTENT):
    slide = prs.slides.add_slide(prs.slide_layouts[layout])
    if title is not None and slide.shapes.title is not None:
        slide.shapes.title.text = title
        set_run_size(slide.shapes.title.text_frame.paragraphs[0], title_size)
    if body_lines:
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = body_lines[0]
        set_run_size(tf.paragraphs[0], body_size)
        for line in body_lines[1:]:
            p = tf.add_paragraph()
            p.text = line
            set_run_size(p, body_size)
    return slide


def add_footer(slide):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(6.95), Inches(8), Inches(0.4))
    tf = box.text_frame
    tf.text = FOOTER
    set_run_size(tf.paragraphs[0], 14)


def make_logo(path, fmt):
    img = Image.new("RGBA", (240, 240), (255, 255, 255, 0))
    d = ImageDraw.Draw(img)
    d.rounded_rectangle([10, 10, 230, 230], radius=40, fill=PRIMARY)
    d.rounded_rectangle([70, 130, 230, 190], radius=20, fill=ACCENT)
    if fmt == "PNG":
        img.save(path, "PNG")
    else:
        img.convert("RGB").save(path, "JPEG")


# ---------------------------------------------------------------- with-skill

def build_with_skill():
    prs = Presentation()
    logo = HERE / "_tmp_logo.png"
    make_logo(logo, "PNG")

    s1 = prs.slides.add_slide(prs.slide_layouts[TITLE_SLIDE])
    s1.shapes.title.text = "Support Ops Proposal for Northwind Logistics"
    set_run_size(s1.shapes.title.text_frame.paragraphs[0], 40)
    if len(s1.placeholders) > 1:
        sub = s1.placeholders[1]
        sub.text_frame.text = "Q3 2026 Review, prepared for Northwind Logistics"
        set_run_size(sub.text_frame.paragraphs[0], 20)
    s1.shapes.add_picture(str(logo), Inches(11.2), Inches(0.4), width=Inches(1.4), height=Inches(1.4))

    s2 = add_slide(prs, title="42% of Q3 Tickets Were Password Resets",
                    body_lines=["1,201 of 2,860 support tickets in Q3 traced to a single resettable flow."])
    add_footer(s2)

    s3 = add_slide(prs, title="Current State vs. Proposed State", body_lines=None)
    left = s3.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(5.7), Inches(4.5))
    ltf = left.text_frame
    ltf.text = "Current"
    set_run_size(ltf.paragraphs[0], 18)
    p = ltf.add_paragraph()
    p.text = "Manual triage across all 11 hubs. Average resolution time: 6.2 hours."
    set_run_size(p, 16)
    right = s3.shapes.add_textbox(Inches(6.9), Inches(1.7), Inches(5.7), Inches(4.5))
    rtf = right.text_frame
    rtf.text = "Proposed"
    set_run_size(rtf.paragraphs[0], 18)
    p = rtf.add_paragraph()
    p.text = "AI-assisted triage, piloted in 3 hubs. Average resolution time: 2.1 hours."
    set_run_size(p, 16)
    add_footer(s3)

    s4 = add_slide(prs, title="Resolution Time Dropped 66% in the Pilot Hubs", body_lines=None)
    chart_data = CategoryChartData()
    chart_data.categories = ["Before pilot", "Pilot hubs"]
    chart_data.add_series("Avg. resolution time (hrs)", (6.2, 2.1))
    s4.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.7), Inches(1.8), Inches(7.5), Inches(4.5), chart_data)
    takeaway = s4.shapes.add_textbox(Inches(8.4), Inches(2.5), Inches(4.2), Inches(3))
    ttf = takeaway.text_frame
    ttf.word_wrap = True
    ttf.text = "Pilot hubs resolved tickets in a third of the time it took the rest of the network."
    set_run_size(ttf.paragraphs[0], 18)
    add_footer(s4)

    s5 = add_slide(prs, title="Rollout Plan", body_lines=[
        "Expand from 3 pilot hubs to all 11 hubs by end of Q1.",
        "Migrate the password-reset flow first; it is 42% of ticket volume.",
        "Train regional leads during the first two weeks of rollout.",
    ])
    add_footer(s5)

    s6 = add_slide(prs, title="Decision Needed by August 15", body_lines=[
        "Approve budget for the full 11-hub rollout.",
        "Confirm the training schedule with regional leads.",
    ])
    add_footer(s6)

    prs.save(HERE / "with-skill.pptx")
    logo.unlink()
    print(f"wrote {HERE / 'with-skill.pptx'}")


# ------------------------------------------------------------- without-skill

BUZZWORD_SLIDES = [
    ("Our Solution Drives Real Value",
     ["We deliver a robust solution that empowers your team.",
      "Our platform is built for today's fast-paced world."]),
    ("Our Solution Drives Real Growth",
     ["Unlock the power of seamless collaboration.",
      "A true game-changer for how support works."]),
    ("Our Solution Drives Real Synergy",
     ["Cutting-edge technology meets world-class service.",
      "We help you leverage synergy across the org."]),
    ("Our Solution Drives Real Impact",
     ["A holistic approach to transformative growth.",
      "Best-in-class innovation, at the end of the day."]),
]

FILLER_SLIDES = [
    ("Why Choose Us", ["We are passionate about excellence.", "Our team is dedicated to your success."]),
    ("Our Approach", ["We believe in strategic alignment.", "Our vision drives our mission forward."]),
    ("Next Steps", ["Let's dive in and get started.", "We're excited to unlock new possibilities."]),
    ("Looking Ahead", ["The future is bright for this partnership.", "Together we can achieve great things."]),
    ("Our Commitment", ["We are committed to your growth journey.", "Excellence is at the core of everything we do."]),
]


def build_without_skill():
    prs = Presentation()
    logo = HERE / "_tmp_logo_bad.jpg"
    make_logo(logo, "JPEG")

    s1 = prs.slides.add_slide(prs.slide_layouts[TITLE_SLIDE])
    s1.shapes.title.text = "Unlocking Seamless Support in a Fast-Paced World!"
    set_run_size(s1.shapes.title.text_frame.paragraphs[0], 40)
    s1.shapes.add_picture(str(logo), Inches(11.2), Inches(0.4), width=Inches(1.4), height=Inches(1.4))

    for title, body in BUZZWORD_SLIDES:
        add_slide(prs, title=title, body_lines=body)

    add_slide(prs, title="Is This the Future of Support?",
              body_lines=["Imagine a world where every ticket resolves itself."])

    add_slide(prs, title="Pricing", body_lines=None)  # title-only

    prs.slides.add_slide(prs.slide_layouts[BLANK])  # empty slide

    s_small = add_slide(prs, title="Implementation Timeline",
                         body_lines=["We roll out in phases across your organization."],
                         body_size=10)  # below the 14pt minimum

    for title, body in FILLER_SLIDES:
        add_slide(prs, title=title, body_lines=body)

    add_slide(prs, title="Thank You!", body_lines=None)

    prs.save(HERE / "without-skill.pptx")
    logo.unlink()
    print(f"wrote {HERE / 'without-skill.pptx'}")


if __name__ == "__main__":
    build_with_skill()
    build_without_skill()
